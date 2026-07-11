#!/usr/bin/env python3
"""
Reclaim — hackathon deck generator.

Builds a 16:9 PowerPoint that mirrors the app's design system (stone-950
dark, amber-500 accent, the five-bar EqMark voice signature, glass cards,
mono tech chips) and leans on two PowerPoint features:

  * Morph transitions on every slide — five persistent 3D bars (the app's
    equalizer mark, extruded with DrawingML scene3d/sp3d and a rotating
    camera) travel, spin and re-equalize between slides.
  * Injected timing XML for fade / rise entrances that echo the app's
    `rise-in` / `materialize` animations.

Run:  python3 presentation/build_deck.py
Out:  presentation/reclaim-hackathon.pptx
"""

import copy
from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

EMU = 914400  # per inch
SLIDE_W, SLIDE_H = 13.333, 7.5

# ---- Reclaim palette (Tailwind stone + amber, straight from globals.css/ui.tsx)
BG      = "0C0A09"  # stone-950 --background
INK     = "FAFAF9"  # --foreground
AMBER   = "F59E0B"  # amber-500 (primary)
AMBER4  = "FBBF24"  # amber-400
AMBER6  = "D97706"  # amber-600
AMBER7  = "B45309"  # amber-700 (extrusion sides)
AMBER1  = "FEF3C7"  # amber-100 (feature titles)
STONE9  = "1C1917"
STONE8  = "292524"
STONE7  = "44403C"
STONE6  = "57534E"
STONE5  = "78716C"
STONE4  = "A8A29E"
STONE3  = "D6D3D1"
STONE2  = "E7E5E4"
ROSE    = "F43F5E"  # rose-500 (recording)
ROSE6   = "E11D48"  # rose-600 (danger)
DOBLUE  = "0080FF"  # DigitalOcean brand blue

SANS = "Segoe UI"
MONO = "Consolas"

NS = ('xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
      'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
      'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
      'xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006" '
      'xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" '
      'xmlns:p159="http://schemas.microsoft.com/office/powerpoint/2015/09/main"')


def frag(xml: str):
    return etree.fromstring(f'<root {NS}>{xml}</root>')[0]


def IN(v: float) -> int:
    return int(round(v * EMU))


def deg(v: float) -> int:
    """Degrees -> positive 60000ths."""
    return int(round(v % 360 * 60000))


# --------------------------------------------------------------------------
# low-level styling helpers (raw DrawingML for alpha, gradients, 3D)
# --------------------------------------------------------------------------

def clear_style(shape):
    el = shape._element
    st = el.find(qn('p:style'))
    if st is not None:
        el.remove(st)


def _clr(hexv, alpha=None):
    a = f'<a:alpha val="{int(alpha*1000)}"/>' if alpha is not None else ''
    return f'<a:srgbClr val="{hexv}">{a}</a:srgbClr>'


def fill_solid(shape, hexv, alpha=None):
    shape._element.spPr.append(frag(f'<a:solidFill>{_clr(hexv, alpha)}</a:solidFill>'))


def fill_grad_v(shape, top, bottom, alpha_top=None, alpha_bottom=None, angle=90):
    """Linear gradient, angle in degrees (90 = top->bottom)."""
    shape._element.spPr.append(frag(
        f'<a:gradFill rotWithShape="1"><a:gsLst>'
        f'<a:gs pos="0">{_clr(top, alpha_top)}</a:gs>'
        f'<a:gs pos="100000">{_clr(bottom, alpha_bottom)}</a:gs>'
        f'</a:gsLst><a:lin ang="{deg(angle)}" scaled="1"/></a:gradFill>'))


def fill_radial(shape, center_hex, center_alpha, edge_hex, edge_alpha):
    shape._element.spPr.append(frag(
        f'<a:gradFill><a:gsLst>'
        f'<a:gs pos="0">{_clr(center_hex, center_alpha)}</a:gs>'
        f'<a:gs pos="100000">{_clr(edge_hex, edge_alpha)}</a:gs>'
        f'</a:gsLst><a:path path="circle">'
        f'<a:fillToRect l="50000" t="50000" r="50000" b="50000"/>'
        f'</a:path></a:gradFill>'))


def fill_none(shape):
    shape._element.spPr.append(frag('<a:noFill/>'))


def line(shape, hexv, w_pt, alpha=None):
    shape._element.spPr.append(frag(
        f'<a:ln w="{int(w_pt*12700)}"><a:solidFill>{_clr(hexv, alpha)}</a:solidFill></a:ln>'))


def line_none(shape):
    shape._element.spPr.append(frag('<a:ln><a:noFill/></a:ln>'))


def shadow(shape, blur=0.18, dist=0.055, alpha=38):
    shape._element.spPr.append(frag(
        f'<a:effectLst><a:outerShdw blurRad="{IN(blur)}" dist="{IN(dist)}" dir="5400000" '
        f'rotWithShape="0"><a:srgbClr val="000000"><a:alpha val="{alpha*1000}"/></a:srgbClr>'
        f'</a:outerShdw></a:effectLst>'))


def make3d(shape, lat=0, lon=0, rev=0, extrusion_in=0.16, contour=AMBER7):
    """Real DrawingML 3D: perspective camera + extruded solid. Morph
    interpolates the camera rotation between slides -> the model spins."""
    spPr = shape._element.spPr
    spPr.append(frag(
        f'<a:scene3d>'
        f'<a:camera prst="perspectiveRelaxedModerately" fov="2700000">'
        f'<a:rot lat="{deg(lat)}" lon="{deg(lon)}" rev="{deg(rev)}"/></a:camera>'
        f'<a:lightRig rig="threePt" dir="t"><a:rot lat="0" lon="0" rev="1200000"/></a:lightRig>'
        f'</a:scene3d>'))
    spPr.append(frag(
        f'<a:sp3d extrusionH="{IN(extrusion_in)}" prstMaterial="matte">'
        f'<a:bevelT w="12700" h="12700"/>'
        f'<a:extrusionClr><a:srgbClr val="{contour}"/></a:extrusionClr>'
        f'</a:sp3d>'))


# --------------------------------------------------------------------------
# shape + text helpers
# --------------------------------------------------------------------------

def box(slide, x, y, w, h, kind=MSO_SHAPE.ROUNDED_RECTANGLE, radius=None, name=None):
    shp = slide.shapes.add_shape(kind, IN(x), IN(y), IN(w), IN(h))
    clear_style(shp)
    if radius is not None and kind == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if name:
        shp.name = name
    shp.shadow.inherit = False
    return shp


def text(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         name=None, wrap=True):
    """paras: list of paragraphs; each a dict(runs=[(txt, opts)], spacing, space_after).
    run opts: size, color, bold, italic, font, spc (tracking, 1/100pt)."""
    tb = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    if name:
        tb.name = name
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get('align', align)
        if para.get('spacing'):
            p.line_spacing = para['spacing']
        if para.get('space_before'):
            p.space_before = Pt(para['space_before'])
        if para.get('space_after') is not None:
            p.space_after = Pt(para['space_after'])
        for txt_, o in para['runs']:
            r = p.add_run()
            r.text = txt_
            f = r.font
            f.name = o.get('font', SANS)
            f.size = Pt(o.get('size', 14))
            f.bold = o.get('bold', False)
            f.italic = o.get('italic', False)
            f.color.rgb = __import__('pptx.dml.color', fromlist=['RGBColor']).RGBColor.from_string(o.get('color', INK))
            if o.get('spc'):
                r._r.rPr.set('spc', str(o['spc']))
    return tb


def P(runs, **kw):
    d = {'runs': runs, 'space_after': kw.pop('space_after', 0)}
    d.update(kw)
    return d


def R(t, **o):
    return (t, o)


# --------------------------------------------------------------------------
# Reclaim components (rebuilt from src/components/ui.tsx + pages)
# --------------------------------------------------------------------------

def bg(slide):
    """stone-950 canvas + faint amber radial glow (morph-persistent)."""
    b = box(slide, 0, 0, SLIDE_W, SLIDE_H, MSO_SHAPE.RECTANGLE, name="!!bg")
    fill_solid(b, BG)
    line_none(b)
    g = box(slide, 3.3, -2.6, 6.8, 6.8, MSO_SHAPE.OVAL, name="!!glow")
    fill_radial(g, AMBER, 9, BG, 0)
    line_none(g)
    return b, g


def eq_bars(slide, cx, baseline, bw, gap, heights, lat=6, lon=-18, rev=0,
            alpha=None, flat=False):
    """The EqMark voice signature as five 3D extruded pills.
    Named !!bar0..4 so Morph tracks them across every slide."""
    total = 5 * bw + 4 * gap
    x0 = cx - total / 2
    shades = [AMBER6, AMBER, AMBER4, AMBER, AMBER6]
    out = []
    for i, h in enumerate(heights):
        x = x0 + i * (bw + gap)
        s = box(slide, x, baseline - h, bw, h, radius=0.5, name=f"!!bar{i}")
        fill_grad_v(s, AMBER4, shades[i], alpha, alpha)
        line_none(s)
        if not flat:
            make3d(s, lat=lat, lon=lon, rev=rev, extrusion_in=max(0.05, bw * 0.5))
        out.append(s)
    return out


def eq_mini(slide, x, y, h=0.16, live_heights=(0.55, 0.85, 1.0, 0.7, 0.45)):
    """Static mini EqMark (as rendered in the app nav)."""
    bw, gap = 0.035, 0.022
    for i, f in enumerate(live_heights):
        s = box(slide, x + i * (bw + gap), y + h * (1 - f), bw, h * f, radius=0.5)
        fill_solid(s, AMBER)
        line_none(s)


def card(slide, x, y, w, h, radius=0.075, fill=STONE9, alpha=90, border=STONE8):
    c = box(slide, x, y, w, h, radius=radius)
    fill_solid(c, fill, alpha)
    line(c, border, 1)
    shadow(c)
    # lit top edge (border-t-stone-700/80 in the app)
    edge = box(slide, x + w * 0.08, y, w * 0.84, 0.016, MSO_SHAPE.RECTANGLE)
    fill_solid(edge, STONE7, 80)
    line_none(edge)
    return c


def chip(slide, x, y, w, label, color=AMBER, mono=True, border=STONE7,
         fillv=STONE9, filla=80, size=10.5, name=None, dot=None, h=0.34):
    c = box(slide, x, y, w, h, radius=0.5, name=name)
    fill_solid(c, fillv, filla)
    line(c, border, 1)
    tx = x
    if dot:
        d = box(slide, x + 0.14, y + h / 2 - 0.045, 0.09, 0.09, MSO_SHAPE.OVAL)
        fill_solid(d, dot)
        line_none(d)
        tx += 0.10
    text(slide, tx, y, w - (tx - x), h,
         [P([R(label, font=MONO if mono else SANS, size=size, color=color, spc=60)])],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return c


def kicker(slide, x, y, label, color=AMBER):
    text(slide, x, y, 8, 0.3,
         [P([R(label, font=MONO, size=12, color=color, bold=True, spc=260)])])


def ghost_num(slide, n):
    text(slide, 10.1, 4.9, 3.4, 2.6,
         [P([R(n, size=170, bold=True, color="171310")])],
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM, wrap=False)


def footer(slide, idx):
    text(slide, 0.7, 7.06, 6, 0.3,
         [P([R("RECLAIM", font=MONO, size=9, color=STONE6, spc=200),
             R("  ·  AI FOR SOCIAL GOOD", font=MONO, size=9, color=STONE6, spc=200)])])
    text(slide, 6.63, 7.06, 6, 0.3,
         [P([R("BUILT ON DIGITALOCEAN  ·  " + idx, font=MONO, size=9, color=STONE6, spc=200)])],
         align=PP_ALIGN.RIGHT)


def btn_primary(slide, x, y, w, label, h=0.52, size=14):
    b = box(slide, x, y, w, h, radius=0.28)
    fill_solid(b, AMBER)
    line_none(b)
    shadow(b, blur=0.12, dist=0.03, alpha=30)
    text(slide, x, y, w, h, [P([R(label, size=size, bold=True, color="1C1917")])],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return b


def btn_secondary(slide, x, y, w, label, h=0.52, size=14):
    b = box(slide, x, y, w, h, radius=0.28)
    fill_solid(b, STONE9, 60)
    line(b, STONE7, 1.2)
    text(slide, x, y, w, h, [P([R(label, size=size, bold=False, color=STONE2)])],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return b


def btn_danger(slide, x, y, w, label, h=0.52, size=14):
    b = box(slide, x, y, w, h, radius=0.28)
    fill_solid(b, ROSE6)
    line_none(b)
    shadow(b, blur=0.12, dist=0.03, alpha=30)
    text(slide, x, y, w, h, [P([R(label, size=size, bold=True, color="FFFFFF")])],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return b


def arrow(slide, x1, y1, x2, y2, color=STONE5, dashed=False, both=False,
          elbow=False, w_pt=1.4, alpha=85):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.ELBOW if elbow else MSO_CONNECTOR.STRAIGHT,
        IN(x1), IN(y1), IN(x2), IN(y2))
    el = conn._element
    st = el.find(qn('p:style'))
    if st is not None:
        el.remove(st)
    dash = '<a:prstDash val="dash"/>' if dashed else ''
    head = '<a:headEnd type="triangle" w="med" len="med"/>' if both else '<a:headEnd type="none"/>'
    el.spPr.append(frag(
        f'<a:ln w="{int(w_pt*12700)}" cap="rnd"><a:solidFill>{_clr(color, alpha)}</a:solidFill>'
        f'{dash}{head}<a:tailEnd type="triangle" w="med" len="med"/></a:ln>'))
    conn.shadow.inherit = False
    return conn


def record_button(slide, cx, cy, r=0.30):
    """The capture button: rose circle + stop square + pulse rings."""
    for rr, al in ((r + 0.22, 12), (r + 0.11, 28)):
        ring = box(slide, cx - rr, cy - rr, rr * 2, rr * 2, MSO_SHAPE.OVAL)
        fill_none(ring)
        line(ring, ROSE, 1.6, al)
    c = box(slide, cx - r, cy - r, r * 2, r * 2, MSO_SHAPE.OVAL)
    fill_solid(c, ROSE6)
    line_none(c)
    shadow(c, blur=0.14, dist=0.02, alpha=45)
    sq = box(slide, cx - 0.085, cy - 0.085, 0.17, 0.17, radius=0.25)
    fill_solid(sq, "FFFFFF")
    line_none(sq)


TILE = 0.78  # architecture icon tile size


def tile(slide, cx, cy, color, label, glyph=None, text_dark=False, img=None,
         img_size=0.44):
    """AWS-diagram-style service tile: colored rounded square + brand icon
    (PNG in presentation/icons/) + two-line label beneath. Returns all
    shapes (for entrance animation)."""
    shapes = []
    t = box(slide, cx - TILE / 2, cy - TILE / 2, TILE, TILE, radius=0.18)
    fill_grad_v(t, color, color)
    line(t, "FFFFFF", 1, 14)
    shadow(t, blur=0.14, dist=0.045, alpha=45)
    shapes.append(t)
    ink = "1C1917" if text_dark else "FFFFFF"
    if img:
        pic = slide.shapes.add_picture(
            f"presentation/icons/{img}.png",
            IN(cx - img_size / 2), IN(cy - img_size / 2), IN(img_size), IN(img_size))
        pic.shadow.inherit = False
        shapes.append(pic)
    elif glyph:
        shapes += glyph(slide, cx, cy, ink)
    lbl = text(slide, cx - 1.05, cy + TILE / 2 + 0.05, 2.1, 0.55,
               [P([R(label, size=10.5, color=STONE3)], spacing=1.05)],
               align=PP_ALIGN.CENTER)
    shapes.append(lbl)
    return shapes


def g_video(s, cx, cy, ink):
    r1 = box(s, cx - 0.20, cy - 0.13, 0.28, 0.26, radius=0.25)
    fill_none(r1); line(r1, ink, 2)
    tr = box(s, cx + 0.10, cy - 0.10, 0.14, 0.20, MSO_SHAPE.ISOSCELES_TRIANGLE)
    tr.rotation = 270
    fill_solid(tr, ink); line_none(tr)
    return [r1, tr]


def g_star(s, cx, cy, ink):
    st = box(s, cx - 0.17, cy - 0.17, 0.34, 0.34, MSO_SHAPE.STAR_4_POINT)
    fill_solid(st, ink); line_none(st)
    return [st]


def g_dots(s, cx, cy, ink):
    out = []
    for dx, dy, r in [(-0.13, -0.11, 0.05), (0.12, -0.07, 0.065),
                      (-0.03, 0.10, 0.05), (0.10, 0.13, 0.04)]:
        d = box(s, cx + dx - r, cy + dy - r, r * 2, r * 2, MSO_SHAPE.OVAL)
        fill_solid(d, ink); line_none(d)
        out.append(d)
    return out


def g_bolt(s, cx, cy, ink):
    b = box(s, cx - 0.13, cy - 0.17, 0.26, 0.34, MSO_SHAPE.LIGHTNING_BOLT)
    fill_solid(b, ink); line_none(b)
    return [b]


def g_bucket(s, cx, cy, ink):
    c = box(s, cx - 0.14, cy - 0.16, 0.28, 0.32, MSO_SHAPE.CAN)
    fill_none(c); line(c, ink, 2)
    return [c]


def g_eq(s, cx, cy, ink):
    out = []
    bw, gap, H = 0.045, 0.028, 0.34
    x0 = cx - (5 * bw + 4 * gap) / 2
    for i, f in enumerate([0.55, 0.85, 1.0, 0.7, 0.45]):
        b = box(s, x0 + i * (bw + gap), cy + H / 2 - H * f, bw, H * f, radius=0.5)
        fill_solid(b, ink); line_none(b)
        out.append(b)
    return out


def g_speaker(s, cx, cy, ink):
    r1 = box(s, cx - 0.17, cy - 0.07, 0.10, 0.14, MSO_SHAPE.RECTANGLE)
    fill_solid(r1, ink); line_none(r1)
    tr = box(s, cx - 0.12, cy - 0.15, 0.22, 0.30, MSO_SHAPE.ISOSCELES_TRIANGLE)
    tr.rotation = 90
    fill_solid(tr, ink); line_none(tr)
    a1 = box(s, cx + 0.13, cy - 0.09, 0.035, 0.18, radius=0.5)
    fill_solid(a1, ink); line_none(a1)
    return [r1, tr, a1]


def g_chip(s, cx, cy, ink):
    outer = box(s, cx - 0.15, cy - 0.15, 0.30, 0.30, radius=0.12)
    fill_none(outer); line(outer, ink, 2)
    inner = box(s, cx - 0.06, cy - 0.06, 0.12, 0.12, MSO_SHAPE.RECTANGLE)
    fill_solid(inner, ink); line_none(inner)
    return [outer, inner]


def g_cloud(s, cx, cy, ink):
    c = box(s, cx - 0.20, cy - 0.13, 0.40, 0.26, MSO_SHAPE.CLOUD)
    fill_none(c); line(c, ink, 2)
    return [c]


def g_lock(s, cx, cy, ink):
    ring = box(s, cx - 0.09, cy - 0.17, 0.18, 0.18, MSO_SHAPE.DONUT)
    fill_solid(ring, ink); line_none(ring)
    body = box(s, cx - 0.13, cy - 0.04, 0.26, 0.20, radius=0.2)
    fill_solid(body, ink); line_none(body)
    return [ring, body]


def g_text(s, cx, cy, ink, label="HF"):
    t = text(s, cx - 0.3, cy - 0.17, 0.6, 0.34,
             [P([R(label, size=15, bold=True, color=ink)])],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return [t]


# --------------------------------------------------------------------------
# morph transition + entrance animation injection
# --------------------------------------------------------------------------

def add_morph(slide, dur=1100):
    slide._element.append(frag(
        f'<mc:AlternateContent>'
        f'<mc:Choice Requires="p159">'
        f'<p:transition spd="slow" p14:dur="{dur}"><p159:morph option="byObject"/></p:transition>'
        f'</mc:Choice>'
        f'<mc:Fallback><p:transition spd="slow"><p:fade/></p:transition></mc:Fallback>'
        f'</mc:AlternateContent>'))


_tid = [100]


def _nid():
    _tid[0] += 1
    return _tid[0]


def _effect(shape_id, delay, dur, rise, first):
    node_type = "afterEffect" if first else "withEffect"
    preset = 42 if rise else 10
    ids = [_nid() for _ in range(4)]
    rise_xml = ''
    if rise:
        rise_xml = (
            f'<p:anim calcmode="lin" valueType="num">'
            f'<p:cBhvr additive="base"><p:cTn id="{ids[3]}" dur="{dur}" fill="hold"/>'
            f'<p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>'
            f'<p:attrNameLst><p:attrName>ppt_y</p:attrName></p:attrNameLst></p:cBhvr>'
            f'<p:tavLst>'
            f'<p:tav tm="0"><p:val><p:strVal val="#ppt_y+0.03"/></p:val></p:tav>'
            f'<p:tav tm="100000"><p:val><p:strVal val="#ppt_y"/></p:val></p:tav>'
            f'</p:tavLst></p:anim>')
    return (
        f'<p:par><p:cTn id="{ids[0]}" presetID="{preset}" presetClass="entr" '
        f'presetSubtype="0" fill="hold" grpId="0" nodeType="{node_type}">'
        f'<p:stCondLst><p:cond delay="{delay}"/></p:stCondLst>'
        f'<p:childTnLst>'
        f'<p:set><p:cBhvr><p:cTn id="{ids[1]}" dur="1" fill="hold">'
        f'<p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
        f'<p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>'
        f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst></p:cBhvr>'
        f'<p:to><p:strVal val="visible"/></p:to></p:set>'
        f'<p:animEffect transition="in" filter="fade">'
        f'<p:cBhvr><p:cTn id="{ids[2]}" dur="{dur}"/>'
        f'<p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl></p:cBhvr></p:animEffect>'
        f'{rise_xml}'
        f'</p:childTnLst></p:cTn></p:par>')


def add_entrances(slide, items):
    """items: list of (shape, delay_ms, rise). Auto-plays after the morph."""
    if not items:
        return
    effects = ''.join(
        _effect(s.shape_id, delay, 550, rise, i == 0)
        for i, (s, delay, rise) in enumerate(items))
    slide._element.append(frag(
        f'<p:timing><p:tnLst><p:par>'
        f'<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>'
        f'<p:seq concurrent="1" nextAc="seek">'
        f'<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>'
        f'<p:par><p:cTn id="3" fill="hold">'
        f'<p:stCondLst><p:cond delay="indefinite"/>'
        f'<p:cond evt="onBegin" delay="0"><p:tn val="2"/></p:cond></p:stCondLst>'
        f'<p:childTnLst><p:par><p:cTn id="4" fill="hold">'
        f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
        f'<p:childTnLst>{effects}</p:childTnLst>'
        f'</p:cTn></p:par></p:childTnLst>'
        f'</p:cTn></p:par>'
        f'</p:childTnLst></p:cTn>'
        f'<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
        f'<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
        f'</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>'))


# --------------------------------------------------------------------------
# deck
# --------------------------------------------------------------------------

prs = Presentation()
prs.slide_width = Emu(IN(SLIDE_W))
prs.slide_height = Emu(IN(SLIDE_H))
BLANK = prs.slide_layouts[6]


def new_slide():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    return s


# ============================== 1 · TITLE ==================================
s = new_slide()
eq_bars(s, cx=6.667, baseline=2.62, bw=0.34, gap=0.21,
        heights=[0.95, 1.55, 1.95, 1.30, 0.85], lat=8, lon=-22, rev=0)

c1 = chip(s, 4.87, 0.5, 3.6, "AI FOR SOCIAL GOOD · HACKATHON", color=STONE3, name="!!hackchip")
t1 = text(s, 1.0, 2.92, 11.333, 1.35,
          [P([R("Reclaim", size=72, bold=True, color=INK)])],
          align=PP_ALIGN.CENTER)
t2 = text(s, 1.0, 4.32, 11.333, 0.75,
          [P([R("Your signs. Your words. ", size=27, color=STONE3),
              R("Your voice.", size=27, bold=True, color=AMBER4)], spacing=1.15)],
          align=PP_ALIGN.CENTER)
t3 = text(s, 2.55, 5.14, 8.25, 0.85,
          [P([R("Sign a phrase to the camera and the room hears it — in your slang, "
                "in a clone of your own voice.", size=15, color=STONE4)], spacing=1.3)],
          align=PP_ALIGN.CENTER)

cx0 = 2.42
c2 = chip(s, cx0, 6.28, 2.55, "BUILT ON DIGITALOCEAN", color="66B2FF", dot=DOBLUE, name="!!dochip")
c3 = chip(s, cx0 + 2.75, 6.28, 1.35, "GEMINI", color=STONE3)
c4 = chip(s, cx0 + 4.30, 6.28, 2.15, "LLAMA 3.3 70B", color=STONE3)
c5 = chip(s, cx0 + 6.65, 6.28, 1.85, "ELEVENLABS", color=STONE3)
add_morph(s)
add_entrances(s, [(c1, 0, False), (t1, 150, True), (t2, 380, True),
                  (t3, 560, True), (c2, 760, False), (c3, 860, False),
                  (c4, 940, False), (c5, 1020, False)])

# ============================== 2 · PROBLEM ================================
s = new_slide()
ghost_num(s, "01")
eq_bars(s, cx=11.85, baseline=1.42, bw=0.13, gap=0.08,
        heights=[0.34, 0.58, 0.74, 0.48, 0.30], lat=4, lon=-40, rev=0, alpha=70)
kicker(s, 0.9, 0.78, "THE PROBLEM")
h = text(s, 0.9, 1.22, 11.6, 2.2,
         [P([R("For millions of Deaf and nonspeaking people,", size=36, bold=True, color=INK)],
            spacing=1.12, space_after=2),
          P([R("every room is a locked door.", size=36, bold=True, color=AMBER4)], spacing=1.12)])
cards_spec = [
    ("70M+ people sign first",
     "Around 70 million Deaf people worldwide use a sign language as their first "
     "language. Most hearing people around them don't."),
    ("TTS speaks — as nobody",
     "Typing into a text-to-speech app is slow, breaks eye contact, and comes out "
     "in a robotic voice that belongs to no one."),
    ("Identity is what's lost",
     "AAC devices hand users a stranger's voice. The words get through — the "
     "person doesn't."),
]
anim = []
for i, (title, body) in enumerate(cards_spec):
    x = 0.9 + i * 3.98
    card(s, x, 3.75, 3.72, 2.5)
    tt = text(s, x + 0.3, 4.05, 3.12, 2.0,
              [P([R(title, size=15.5, bold=True, color=AMBER1)], space_after=8),
               P([R(body, size=12.5, color=STONE4)], spacing=1.32)])
    anim.append((tt, 300 + i * 220, True))
footer(s, "02 / 11")
add_morph(s)
add_entrances(s, [(h, 0, True)] + anim)

# ============================== 3 · THE IDEA ===============================
s = new_slide()
ghost_num(s, "02")
eq_bars(s, cx=1.72, baseline=6.55, bw=0.22, gap=0.135,
        heights=[0.68, 1.12, 1.46, 0.92, 0.62], lat=2, lon=28, rev=0)
kicker(s, 0.9, 0.78, "THE IDEA")
h = text(s, 0.9, 1.22, 7.6, 2.2,
         [P([R("Sign it.", size=44, bold=True, color=INK)], spacing=1.08, space_after=4),
          P([R("The room hears ", size=44, bold=True, color=INK),
             R("you", size=44, bold=True, color=AMBER4),
             R(".", size=44, bold=True, color=INK)], spacing=1.08)])
sub = text(s, 0.9, 3.55, 6.6, 1.5,
           [P([R("Record a clip of yourself signing. Reclaim translates it, lets you "
                 "fix every word, rewrites it the way you actually talk, and says it "
                 "out loud — in a voice cloned from you.", size=15, color=STONE4)],
              spacing=1.4)])
idea_cards = [
    ("YOUR WORDS", "A style model trained on how you actually text — your slang, your punctuation, your energy."),
    ("YOUR VOICE", "Cloned from 1–3 minutes of speech. A library of named voices; switch anytime."),
    ("YOUR CONTROL", "Human in the loop: nothing is ever spoken without your review."),
]
anim = []
for i, (title, body) in enumerate(idea_cards):
    y = 1.35 + i * 1.72
    card(s, 8.35, y, 4.1, 1.5)
    tt = text(s, 8.65, y + 0.24, 3.5, 1.1,
              [P([R(title, font=MONO, size=11, bold=True, color=AMBER, spc=160)], space_after=6),
               P([R(body, size=12, color=STONE4)], spacing=1.28)])
    anim.append((tt, 350 + i * 220, True))
footer(s, "03 / 11")
add_morph(s)
add_entrances(s, [(h, 0, True), (sub, 200, True)] + anim)

# ============================== 4 · PIPELINE ===============================
s = new_slide()
ghost_num(s, "03")
eq_bars(s, cx=12.10, baseline=1.30, bw=0.13, gap=0.08,
        heights=[0.30, 0.52, 0.68, 0.44, 0.28], lat=0, lon=-60, rev=0, alpha=70)
kicker(s, 0.9, 0.62, "THE PIPELINE")
h = text(s, 0.9, 1.02, 11, 0.75,
         [P([R("From sign to speech in four steps", size=30, bold=True, color=INK)])])
sub = text(s, 0.9, 1.62, 11, 0.4,
           [P([R("Every stage is a different model doing what it's best at — stitched into one tap.",
                 size=13.5, color=STONE4)])])
pipe = [
    ("1", "Sign it",
     "Gemini ingests the recorded clip as native video — signing is motion, so "
     "real video beats sampled frames. Holding up 1–5 fingers fires a quick phrase.",
     "GEMINI · NATIVE VIDEO INPUT"),
    ("2", "Check it",
     "The translation comes back as editable text. Nothing is spoken without "
     "your review — you stay in control of every word.",
     "HUMAN IN THE LOOP"),
    ("3", "Make it yours",
     "An LLM rewrites the sentence the way you text — a distilled style card "
     "plus your most similar past messages, retrieved semantically.",
     "LLAMA 3.3 70B · DO GRADIENT · GTE-LARGE"),
    ("4", "Say it out loud",
     "Spoken in your cloned voice from your voice library. Repeated phrases "
     "return instantly from cache without spending TTS credits.",
     "ELEVENLABS FLASH V2.5 · DO MANAGED VALKEY"),
]
anim = []
for i, (n, title, body, tech) in enumerate(pipe):
    x = 0.9 + (i % 2) * 5.87
    y = 2.28 + (i // 2) * 2.32
    card(s, x, y, 5.62, 2.12)
    num = box(s, x + 0.28, y + 0.26, 0.42, 0.42, MSO_SHAPE.OVAL)
    fill_solid(num, AMBER)
    line_none(num)
    text(s, x + 0.28, y + 0.26, 0.42, 0.42,
         [P([R(n, size=15, bold=True, color="1C1917")])],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tt = text(s, x + 0.88, y + 0.30, 4.5, 0.4,
              [P([R(title, size=16.5, bold=True, color=INK)])])
    bb = text(s, x + 0.28, y + 0.82, 5.05, 1.0,
              [P([R(body, size=11.5, color=STONE4)], spacing=1.25)])
    tc = text(s, x + 0.28, y + 1.74, 5.05, 0.3,
              [P([R(tech, font=MONO, size=9.5, color=AMBER, spc=80)])])
    anim += [(tt, 250 + i * 200, True), (bb, 330 + i * 200, True), (tc, 410 + i * 200, False)]
footer(s, "04 / 11")
add_morph(s)
add_entrances(s, [(h, 0, True), (sub, 120, True)] + anim)

# ============================ 5 · SPEAK SCREEN =============================
s = new_slide()
ghost_num(s, "04")
eq_bars(s, cx=12.28, baseline=6.72, bw=0.11, gap=0.07,
        heights=[0.26, 0.44, 0.58, 0.38, 0.24], lat=0, lon=45, rev=0, alpha=70)
kicker(s, 4.75, 0.78, "THE APP — LIVE")
h = text(s, 4.75, 1.20, 8.0, 0.75,
         [P([R("The Speak screen", size=30, bold=True, color=INK)])])
sub = text(s, 4.75, 1.80, 7.6, 0.5,
           [P([R("Every element on this slide is a component from the running app.",
                 size=13.5, color=STONE4)])])

# ---- phone mockup (Speak page, recording state)
px, py, pw, ph = 1.05, 0.95, 2.95, 5.72
phone = box(s, px, py, pw, ph, radius=0.12)
fill_solid(phone, "141110")
line(phone, STONE7, 1.6)
shadow(phone, blur=0.28, dist=0.08, alpha=55)
# chrome nav (translucent chrome bar + EqMark + wordmark)
nav = box(s, px + 0.12, py + 0.13, pw - 0.24, 0.44, radius=0.18)
fill_solid(nav, BG, 75)
line(nav, "FFFFFF", 0.75, 6)
eq_mini(s, px + 0.30, py + 0.27, h=0.17)
text(s, px + 0.62, py + 0.13, 1.6, 0.44,
     [P([R("Reclaim", size=13, bold=True, color=INK)])], anchor=MSO_ANCHOR.MIDDLE)
# camera viewport
vp = box(s, px + 0.12, py + 0.68, pw - 0.24, 2.42, radius=0.06)
fill_solid(vp, "0A0908")
line(vp, STONE8, 1)
for bx_, by_, bw_, bh_ in [  # corner brackets
    (px+0.26, py+0.82, 0.22, 0.02), (px+0.26, py+0.82, 0.02, 0.22),
    (px+pw-0.48, py+0.82, 0.22, 0.02), (px+pw-0.28, py+0.82, 0.02, 0.22),
    (px+0.26, py+2.88, 0.22, 0.02), (px+0.26, py+2.68, 0.02, 0.22),
    (px+pw-0.48, py+2.88, 0.22, 0.02), (px+pw-0.28, py+2.68, 0.02, 0.22),
]:
    br = box(s, bx_, by_, bw_, bh_, MSO_SHAPE.RECTANGLE)
    fill_solid(br, STONE5, 70)
    line_none(br)
# REC chip
rec = box(s, px + 0.30, py + 0.90, 0.86, 0.30, radius=0.5)
fill_solid(rec, "000000", 60)
line_none(rec)
rd = box(s, px + 0.40, py + 1.005, 0.09, 0.09, MSO_SHAPE.OVAL)
fill_solid(rd, ROSE)
line_none(rd)
text(s, px + 0.53, py + 0.90, 0.6, 0.30,
     [P([R("0:07", font=MONO, size=10, color="FFFFFF")])], anchor=MSO_ANCHOR.MIDDLE)
text(s, px + 0.12, py + 1.95, pw - 0.24, 0.4,
     [P([R("signing…", size=12, italic=True, color=STONE5)])],
     align=PP_ALIGN.CENTER)
# record button with pulse rings
record_button(s, px + pw / 2, py + 3.58, r=0.27)
# translation card
tcard = box(s, px + 0.12, py + 4.12, pw - 0.24, 0.95, radius=0.10)
fill_solid(tcard, STONE9, 90)
line(tcard, STONE8, 1)
text(s, px + 0.28, py + 4.24, pw - 0.56, 0.75,
     [P([R("TRANSLATION — EDIT ANYTHING", font=MONO, size=8, color=AMBER, spc=80)], space_after=4),
      P([R("coffee run — want anything?", size=12.5, color=INK)])])
# buttons
btn_secondary(s, px + 0.12, py + 5.22, 1.28, "Re-sign", h=0.40, size=11)
btn_primary(s, px + 1.51, py + 5.22, 1.32, "Say it", h=0.40, size=11)

# ---- annotations
notes = [
    ("Recording pulse", "The rose capture ring pulses while the camera records your signing — 20 s max, auto-stops."),
    ("Human in the loop", "Gemini's translation lands in an editable card. Fix anything before it's ever spoken."),
    ("One amber tap", "“Say it” restyles the sentence and speaks it in your cloned voice — the room hears you."),
    ("Quick phrases", "Hold up 1–5 fingers as the clip starts: deterministic phrases like “I want a coffee.”"),
]
anim = []
for i, (title, body) in enumerate(notes):
    y = 2.55 + i * 1.08
    bullet = box(s, 4.78, y + 0.09, 0.12, 0.12, radius=0.35)
    fill_solid(bullet, AMBER)
    line_none(bullet)
    tt = text(s, 5.05, y, 7.4, 1.0,
              [P([R(title + " — ", size=13.5, bold=True, color=AMBER1),
                  R(body, size=13, color=STONE4)], spacing=1.25)])
    anim.append((tt, 350 + i * 200, True))
footer(s, "05 / 11")
add_morph(s)
add_entrances(s, [(h, 0, True), (sub, 120, True)] + anim)

# ========================== 6 · CONVERSATION MODE ==========================
s = new_slide()
ghost_num(s, "05")
eq_bars(s, cx=1.35, baseline=6.72, bw=0.11, gap=0.07,
        heights=[0.26, 0.44, 0.58, 0.38, 0.24], lat=0, lon=70, rev=0, alpha=70)
kicker(s, 0.9, 0.78, "CONVERSATION MODE")
h = text(s, 0.9, 1.20, 6.6, 1.6,
         [P([R("A real\nback-and-forth.", size=36, bold=True, color=INK)], spacing=1.1)])
sub = text(s, 0.9, 2.75, 5.9, 2.2,
           [P([R("The other person just talks. Gemini transcribes them; Llama 3.3 on "
                 "DigitalOcean Gradient drafts three replies in your style. Tap one — "
                 "it's spoken in your voice.", size=14.5, color=STONE4)], spacing=1.4)])
flow = text(s, 0.9, 4.55, 5.9, 0.4,
            [P([R("HEAR → DRAFT ×3 → TAP → SPEAK", font=MONO, size=12, color=AMBER, spc=120)])])

# partner bubble
pb = box(s, 7.15, 1.30, 5.25, 0.92, radius=0.16)
fill_solid(pb, STONE8, 95)
line(pb, STONE7, 1)
text(s, 7.45, 1.42, 4.7, 0.7,
     [P([R("PARTNER · TRANSCRIBED BY GEMINI", font=MONO, size=8.5, color=STONE5, spc=80)], space_after=4),
      P([R("“What do you want to do for lunch?”", size=13.5, color=STONE2)])])
# reply chips (middle one active/playing)
replies = [
    ("tbh i could demolish a burrito rn", False),
    ("wherever u want, i'm easy", True),
    ("gotta be quick — i'm slammed today", False),
]
anim = []
lbl = text(s, 7.15, 2.42, 5.25, 0.35,
           [P([R("YOUR REPLIES · DRAFTED IN YOUR STYLE", font=MONO, size=8.5, color=AMBER, spc=80)])])
for i, (t, active) in enumerate(replies):
    y = 2.85 + i * 0.98
    rb = box(s, 7.15, y, 5.25, 0.78, radius=0.5)
    if active:
        fill_solid(rb, AMBER, 14)
        line(rb, AMBER, 1.4)
    else:
        fill_solid(rb, STONE9, 85)
        line(rb, STONE7, 1)
    text(s, 7.50, y, 4.2, 0.78,
         [P([R(t, size=13, color=INK if active else STONE3)])], anchor=MSO_ANCHOR.MIDDLE)
    if active:
        eq_mini(s, 11.72, y + 0.26, h=0.26)
    anim.append((rb, 400 + i * 220, True))
cap = text(s, 7.15, 5.85, 5.25, 0.4,
           [P([R("▲ playing in your cloned voice", font=MONO, size=10, color=AMBER, spc=60)])],
           align=PP_ALIGN.CENTER)
footer(s, "06 / 11")
add_morph(s)
add_entrances(s, [(h, 0, True), (sub, 150, True), (flow, 280, False),
                  (pb, 380, True), (lbl, 500, False)] + anim + [(cap, 1150, False)])

# ========================== 7 · DIGITALOCEAN ===============================
s = new_slide()
ghost_num(s, "06")
eq_bars(s, cx=11.55, baseline=2.05, bw=0.16, gap=0.10,
        heights=[0.40, 0.66, 0.86, 0.56, 0.36], lat=6, lon=-30, rev=8)
kicker(s, 0.9, 0.62, "KEY TECHNOLOGY", color="66B2FF")
h = text(s, 0.9, 1.02, 10.5, 0.75,
         [P([R("DigitalOcean is the backbone", size=30, bold=True, color=INK)])])
sub = text(s, 0.9, 1.62, 9.6, 0.55,
           [P([R("One tap on “Say it” fans out across these services — and degrades "
                 "gracefully at every rung. The flow never blocks.", size=13.5, color=STONE4)],
              spacing=1.3)])
stack = [
    ("GRADIENT INFERENCE", "Llama 3.3 70B — serverless — rewrites every sentence in your own style"),
    ("GRADIENT EMBEDDINGS", "GTE-Large embeds your corpus; the closest past texts become few-shot examples"),
    ("MANAGED VALKEY", "Audio cached per (voice, sentence) — instant repeats, zero wasted TTS credits"),
    ("SPACES", "S3-compatible object storage — private per-user prefix, PII redacted before write"),
    ("APP PLATFORM", "Auto-deployed from GitHub — the live demo is running on it right now"),
    ("GPU DROPLET", "Self-hosted F5-TTS voice cloning — first-party voice biometrics at scale"),
    ("QWEN3 TTS · GRADIENT", "Neutral-voice fallback when a user has no clone yet — nobody is left silent"),
]
tbl = box(s, 0.9, 2.42, 11.53, len(stack) * 0.60 + 0.08, radius=0.045)
fill_solid(tbl, STONE9, 45)
line(tbl, STONE8, 1)
anim = []
for i, (name, desc) in enumerate(stack):
    y = 2.46 + i * 0.60
    if i % 2 == 0:
        row = box(s, 0.94, y, 11.45, 0.60, MSO_SHAPE.RECTANGLE)
        fill_solid(row, STONE9, 75)
        line_none(row)
    mark = box(s, 1.14, y + 0.235, 0.13, 0.13, radius=0.4)
    fill_solid(mark, DOBLUE)
    line_none(mark)
    nm = text(s, 1.45, y, 3.05, 0.60,
              [P([R(name, font=MONO, size=10.5, bold=True, color="66B2FF", spc=60)])],
              anchor=MSO_ANCHOR.MIDDLE)
    ds = text(s, 4.6, y, 7.7, 0.60,
              [P([R(desc, size=12, color=STONE3)])], anchor=MSO_ANCHOR.MIDDLE)
    anim += [(nm, 250 + i * 130, False), (ds, 310 + i * 130, False)]
footer(s, "07 / 11")
add_morph(s)
add_entrances(s, [(h, 0, True), (sub, 120, True), (tbl, 200, True)] + anim)

# ========================== 8 · ARCHITECTURE ===============================
# Service-diagram slide modeled on the classic AWS architecture layout:
# icon tiles, orthogonal arrows, client node left, caption top-left.
s = new_slide()
ghost_num(s, "07")
eq_bars(s, cx=12.55, baseline=0.98, bw=0.10, gap=0.065,
        heights=[0.22, 0.38, 0.50, 0.32, 0.20], lat=0, lon=-75, rev=0, alpha=60)
kicker(s, 0.6, 0.42, "ARCHITECTURE")
hh = text(s, 0.6, 0.80, 7.5, 0.5,
          [P([R("Sign-to-speech serverless backend", size=20, bold=True, color=INK)])])

VIOLET = "8B5CF6"; AWSRED = "DD344C"; EC2ORG = "ED7100"; HFYELL = "FFD21E"; ELDARK = "26211E"

nodes = {}
nodes['auth'] = tile(s, 1.35, 2.60, AWSRED, "AWS Amplify\nAuth", img="awsamplify", img_size=0.62)
# client node — outlined box, as in the reference diagram
cbx, cby, cbw, cbh = 2.225, 3.45, 1.05, 1.40
cb = box(s, cbx, cby, cbw, cbh, MSO_SHAPE.RECTANGLE)
fill_solid(cb, BG, 40); line(cb, STONE3, 1.75)
mon = box(s, 2.53, 3.72, 0.44, 0.34, radius=0.12)
fill_none(mon); line(mon, STONE2, 2)
stand = box(s, 2.70, 4.06, 0.10, 0.09, MSO_SHAPE.RECTANGLE)
fill_solid(stand, STONE2); line_none(stand)
clbl = text(s, cbx, 4.32, cbw, 0.35, [P([R("Client", size=11.5, color=STONE2)])],
            align=PP_ALIGN.CENTER)
csub = text(s, cbx - 0.4, 4.90, cbw + 0.8, 0.3,
            [P([R("Web · Expo mobile", size=9.5, color=STONE5)])], align=PP_ALIGN.CENTER)
nodes['client'] = [cb, mon, stand, clbl, csub]
nodes['api'] = tile(s, 4.70, 4.15, DOBLUE, "Next.js API\nDO App Platform", img="digitalocean")
nodes['gemini'] = tile(s, 4.70, 1.70, VIOLET, "Gemini\nsign recognition", img="googlegemini")
nodes['embed'] = tile(s, 7.00, 1.70, DOBLUE, "Gradient Embeddings\nGTE-Large", img="digitalocean")
nodes['spaces'] = tile(s, 9.40, 1.70, DOBLUE, "DO Spaces\ncorpus · PII-redacted", img="digitalocean")
nodes['infer'] = tile(s, 7.00, 4.15, DOBLUE, "Gradient Inference\nLlama 3.3 70B", img="meta")
nodes['eleven'] = tile(s, 9.40, 4.15, ELDARK, "ElevenLabs\nclone + Flash TTS", img="elevenlabs")
nodes['valkey'] = tile(s, 11.70, 4.15, DOBLUE, "Managed Valkey\nTTS cache", img="digitalocean")
nodes['f5'] = tile(s, 2.75, 6.35, DOBLUE, "DO GPU Droplet\nF5-TTS self-hosted", img="digitalocean")
nodes['ec2'] = tile(s, 5.15, 6.35, EC2ORG, "EC2 GPU\nLoRA fine-tune (TRL)", img="amazonec2", img_size=0.62)
nodes['hf'] = tile(s, 7.55, 6.35, HFYELL, "Hugging Face\ndedicated endpoint", img="huggingface")
nodes['qwen'] = tile(s, 9.95, 6.35, DOBLUE, "Qwen3 TTS\nneutral fallback", img="qwen")

HT = TILE / 2
arrows = [
    arrow(s, cbx + cbw / 2, cby, 1.35 + HT, 2.60, elbow=True),            # client -> auth
    arrow(s, cbx + cbw, 4.15, 4.70 - HT, 4.15, both=True),                # client <-> api
    arrow(s, 4.70, 4.15 - HT, 4.70, 1.70 + HT),                           # api -> gemini
    arrow(s, 4.70 + HT, 4.15, 7.00 - HT, 4.15),                           # api -> inference
    arrow(s, 7.00 + HT, 4.15, 9.40 - HT, 4.15),                           # inference -> elevenlabs
    arrow(s, 9.40 + HT, 4.15, 11.70 - HT, 4.15, both=True),               # elevenlabs <-> valkey
    arrow(s, 9.40 - HT, 1.70, 7.00 + HT, 1.70),                           # spaces -> embeddings
    arrow(s, 7.00, 1.70 + HT, 7.00, 4.15 - HT),                           # embeddings -> inference
    arrow(s, 4.70, 4.15 + HT, 2.75 + HT, 6.35, elbow=True, dashed=True),  # api -> F5 droplet
    arrow(s, 5.15 + HT, 6.35, 7.55 - HT, 6.35),                           # ec2 -> hugging face
    arrow(s, 7.55, 6.35 - HT, 7.00, 4.15 + HT, elbow=True, dashed=True),  # hf -> inference
    arrow(s, 9.40, 4.15 + HT, 9.95, 6.35 - HT, elbow=True, dashed=True),  # elevenlabs -> qwen fallback
]
order = ['client', 'auth', 'api', 'gemini', 'infer', 'embed', 'spaces',
         'eleven', 'valkey', 'qwen', 'f5', 'ec2', 'hf']
anim = [(hh, 0, True)]
for i, key in enumerate(order):
    for shp in nodes[key]:
        anim.append((shp, 200 + i * 150, False))
for i, a in enumerate(arrows):
    anim.append((a, 300 + i * 120, False))
add_morph(s)
add_entrances(s, anim)

# ========================== 9 · CONSENT & PRIVACY ==========================
s = new_slide()
ghost_num(s, "08")
eq_bars(s, cx=11.95, baseline=1.42, bw=0.13, gap=0.08,
        heights=[0.32, 0.55, 0.72, 0.46, 0.30], lat=0, lon=-15, rev=0, alpha=70)
kicker(s, 0.9, 0.78, "SOCIAL GOOD INCLUDES PRIVACY")
h = text(s, 0.9, 1.22, 11, 0.8,
         [P([R("Voiceprints are biometric data.\n", size=30, bold=True, color=INK),
             R("We treat them like it.", size=30, bold=True, color=AMBER4)], spacing=1.15)])
priv = [
    ("Hard consent gate",
     "Voice upload and cloning are blocked until an explicit, written consent "
     "record is stored. GDPR Art. 9 · Illinois BIPA · Texas CUBI."),
    ("PII redaction",
     "URLs, emails and phone numbers are stripped from your text corpus before "
     "it ever touches storage."),
    ("One-tap erasure",
     "Settings → Delete my data removes the Spaces objects, the consent record, "
     "and every ElevenLabs clone — everything, everywhere."),
]
anim = []
for i, (title, body) in enumerate(priv):
    x = 0.9 + i * 3.98
    card(s, x, 3.05, 3.72, 2.15)
    tt = text(s, x + 0.3, 3.35, 3.12, 1.7,
              [P([R(title, size=15.5, bold=True, color=AMBER1)], space_after=8),
               P([R(body, size=12, color=STONE4)], spacing=1.3)])
    anim.append((tt, 300 + i * 220, True))
# the app's actual danger button
db = btn_danger(s, 5.24, 5.70, 2.85, "Delete my data", h=0.55)
cap = text(s, 2.9, 6.38, 7.53, 0.35,
           [P([R("the same rose button that ships in Settings — revocation is a feature, not a form",
                 font=MONO, size=10, color=STONE5, spc=40)])], align=PP_ALIGN.CENTER)
footer(s, "09 / 11")
add_morph(s)
add_entrances(s, [(h, 0, True)] + anim + [(db, 950, True), (cap, 1100, False)])

# ========================= 10 · THE LADDER =================================
s = new_slide()
ghost_num(s, "09")
# bars echo the ladder: ascending heights
eq_bars(s, cx=11.35, baseline=6.30, bw=0.20, gap=0.13,
        heights=[0.45, 0.80, 1.15, 1.50, 1.85], lat=4, lon=22, rev=0)
kicker(s, 0.9, 0.78, "BEYOND THE DEMO")
h = text(s, 0.9, 1.22, 10.5, 0.8,
         [P([R("A ladder to a model that's yours alone", size=30, bold=True, color=INK)])])
sub = text(s, 0.9, 1.85, 9.2, 0.5,
           [P([R("Personal style isn't one trick — it's three tiers, and the top tier is real training.",
                 size=13.5, color=STONE4)])])
tiers = [
    ("TIER 1 · ALWAYS ON", "Few-shot retrieval",
     "Your corpus is embedded once (GTE-Large on DO). At speak time, your most "
     "similar past messages become the examples."),
    ("TIER 2 · PER STYLE", "Distilled style card",
     "An LLM writes a profile of how you text — slang, punctuation, energy — "
     "injected into every prompt. Swappable: “Casual”, “Business”…"),
    ("TIER 3 · PER USER", "Your own fine-tuned model",
     "A LoRA of Qwen2.5-1.5B trained on your messages (TRL, EC2 GPU), pushed to "
     "Hugging Face and served from a dedicated endpoint — a model of one."),
]
anim = []
for i, (tag, title, body) in enumerate(tiers):
    x = 0.9 + i * 3.35
    y = 4.55 - i * 0.85          # each tier steps up
    w = 3.12
    card(s, x, y, w, 2.05 + i * 0.0)
    tt = text(s, x + 0.26, y + 0.26, w - 0.52, 1.7,
              [P([R(tag, font=MONO, size=9.5, bold=True, color=AMBER, spc=100)], space_after=5),
               P([R(title, size=15, bold=True, color=AMBER1)], space_after=6),
               P([R(body, size=11, color=STONE4)], spacing=1.25)])
    anim.append((tt, 300 + i * 260, True))
mob = chip(s, 4.35, 6.38, 4.6, "+ EXPO MOBILE APP — SAME IDENTITY, POCKET-SIZED",
           color=STONE3, size=9.5)
footer(s, "10 / 11")
add_morph(s)
add_entrances(s, [(h, 0, True), (sub, 130, True)] + anim + [(mob, 1150, False)])

# ============================= 11 · CLOSE ==================================
s = new_slide()
eq_bars(s, cx=6.667, baseline=3.30, bw=0.38, gap=0.24,
        heights=[1.55, 2.15, 2.60, 1.95, 1.40], lat=10, lon=18, rev=0)
h = text(s, 1.0, 3.62, 11.333, 1.2,
         [P([R("Five minutes of setup.", size=32, bold=True, color=INK)],
            spacing=1.12, space_after=2),
          P([R("A voice that's yours for good.", size=32, bold=True, color=AMBER4)],
            spacing=1.12)], align=PP_ALIGN.CENTER)
url = chip(s, 4.42, 5.35, 4.5, "reclaim-2p92q.ondigitalocean.app",
           color="1C1917", mono=True, border=AMBER, fillv=AMBER, filla=100,
           size=12, h=0.46)
sub = text(s, 2.5, 6.12, 8.33, 0.5,
           [P([R("Reclaim — AI for social good, built on DigitalOcean. Thank you.",
                 size=14, color=STONE4)])], align=PP_ALIGN.CENTER)
c2 = chip(s, 4.87, 0.5, 3.6, "AI FOR SOCIAL GOOD · HACKATHON", color=STONE3, name="!!hackchip")
add_morph(s)
add_entrances(s, [(h, 0, True), (url, 300, True), (sub, 500, False)])

# --------------------------------------------------------------------------
out = "presentation/reclaim-hackathon.pptx"
prs.save(out)
print(f"wrote {out} with {len(prs.slides._sldIdLst)} slides")
